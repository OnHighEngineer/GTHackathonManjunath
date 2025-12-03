"""
Report Generator Service
Creates PDF and PowerPoint reports from data and insights
"""

import os
from typing import Dict, Any, List
from datetime import datetime
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from config import settings


class ReportGenerator:
    def __init__(self):
        self.output_folder = settings.OUTPUT_FOLDER
        os.makedirs(self.output_folder, exist_ok=True)
        
        # Brand colors
        self.primary_color = colors.HexColor("#2563eb")
        self.secondary_color = colors.HexColor("#1e40af")
        self.accent_color = colors.HexColor("#3b82f6")
        self.dark_color = colors.HexColor("#1f2937")
        self.light_color = colors.HexColor("#f3f4f6")
    
    def generate_report(
        self,
        data: Dict[str, Any],
        insights: Dict[str, Any],
        charts: Dict[str, str],
        config: Dict[str, Any],
        job_id: str
    ) -> str:
        """Generate report in the specified format"""
        report_type = config.get("report_type", "pdf")
        
        if report_type == "pdf":
            return self._generate_pdf(data, insights, charts, config, job_id)
        elif report_type == "pptx":
            return self._generate_pptx(data, insights, charts, config, job_id)
        else:
            raise ValueError(f"Unsupported report type: {report_type}")
    
    def _generate_pdf(
        self,
        data: Dict[str, Any],
        insights: Dict[str, Any],
        charts: Dict[str, str],
        config: Dict[str, Any],
        job_id: str
    ) -> str:
        """Generate a professional PDF report"""
        output_path = os.path.join(self.output_folder, f"{job_id}.pdf")
        
        doc = SimpleDocTemplate(
            output_path,
            pagesize=letter,
            rightMargin=0.75*inch,
            leftMargin=0.75*inch,
            topMargin=0.75*inch,
            bottomMargin=0.75*inch
        )
        
        # Styles
        styles = getSampleStyleSheet()
        
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=28,
            textColor=self.primary_color,
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=self.secondary_color,
            spaceBefore=20,
            spaceAfter=10
        )
        
        subheading_style = ParagraphStyle(
            'CustomSubheading',
            parent=styles['Heading3'],
            fontSize=12,
            textColor=self.dark_color,
            spaceBefore=15,
            spaceAfter=8
        )
        
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['Normal'],
            fontSize=10,
            textColor=self.dark_color,
            alignment=TA_JUSTIFY,
            spaceAfter=8,
            leading=14
        )
        
        bullet_style = ParagraphStyle(
            'CustomBullet',
            parent=styles['Normal'],
            fontSize=10,
            textColor=self.dark_color,
            leftIndent=20,
            spaceAfter=5,
            bulletIndent=10
        )
        
        # Build document content
        story = []
        
        # Title Page
        story.append(Spacer(1, 2*inch))
        story.append(Paragraph(config.get("title", "Performance Report"), title_style))
        story.append(Spacer(1, 0.3*inch))
        story.append(Paragraph(config.get("company_name", "Company"), 
                              ParagraphStyle('Subtitle', parent=styles['Normal'], 
                                           fontSize=14, textColor=self.accent_color, 
                                           alignment=TA_CENTER)))
        story.append(Spacer(1, 0.5*inch))
        story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}", 
                              ParagraphStyle('Date', parent=styles['Normal'], 
                                           fontSize=11, alignment=TA_CENTER)))
        story.append(PageBreak())
        
        # Executive Summary
        if config.get("include_summary", True) and insights.get("executive_summary"):
            story.append(Paragraph("Executive Summary", heading_style))
            story.append(Paragraph(insights["executive_summary"], body_style))
            story.append(Spacer(1, 0.3*inch))
        
        # Key Findings
        if insights.get("key_findings"):
            story.append(Paragraph("Key Findings", heading_style))
            for finding in insights["key_findings"]:
                story.append(Paragraph(f"• {finding}", bullet_style))
            story.append(Spacer(1, 0.3*inch))
        
        # Performance Highlights
        if insights.get("performance_highlights"):
            story.append(Paragraph("Performance Highlights", heading_style))
            
            highlights = insights["performance_highlights"]
            if highlights.get("top_performers"):
                story.append(Paragraph("Top Performers", subheading_style))
                for item in highlights["top_performers"]:
                    story.append(Paragraph(f"✓ {item}", bullet_style))
            
            if highlights.get("areas_of_concern"):
                story.append(Paragraph("Areas of Concern", subheading_style))
                for item in highlights["areas_of_concern"]:
                    story.append(Paragraph(f"⚠ {item}", bullet_style))
            
            story.append(Spacer(1, 0.3*inch))
        
        # Charts Section
        if config.get("include_charts", True) and charts:
            story.append(Paragraph("Visual Analytics", heading_style))
            
            for chart_name, chart_path in charts.items():
                if os.path.exists(chart_path):
                    try:
                        story.append(Spacer(1, 0.2*inch))
                        story.append(Paragraph(chart_name.replace("_", " ").title(), subheading_style))
                        img = Image(chart_path, width=6*inch, height=4*inch)
                        story.append(img)
                        story.append(Spacer(1, 0.2*inch))
                    except Exception as e:
                        print(f"Error adding chart {chart_name}: {e}")
        
        # Trends
        if insights.get("trends"):
            story.append(PageBreak())
            story.append(Paragraph("Trends & Patterns", heading_style))
            for trend in insights["trends"]:
                story.append(Paragraph(f"📈 {trend}", bullet_style))
            story.append(Spacer(1, 0.3*inch))
        
        # Recommendations
        if config.get("include_recommendations", True) and insights.get("recommendations"):
            story.append(Paragraph("Strategic Recommendations", heading_style))
            for i, rec in enumerate(insights["recommendations"], 1):
                story.append(Paragraph(f"{i}. {rec}", bullet_style))
            story.append(Spacer(1, 0.3*inch))
        
        # Risk Factors
        if insights.get("risk_factors"):
            story.append(Paragraph("Risk Factors", heading_style))
            for risk in insights["risk_factors"]:
                story.append(Paragraph(f"⚡ {risk}", bullet_style))
            story.append(Spacer(1, 0.3*inch))
        
        # Opportunities
        if insights.get("opportunities"):
            story.append(Paragraph("Growth Opportunities", heading_style))
            for opp in insights["opportunities"]:
                story.append(Paragraph(f"💡 {opp}", bullet_style))
        
        # Data Summary Table
        story.append(PageBreak())
        story.append(Paragraph("Data Summary", heading_style))
        
        metadata = data.get("metadata", {})
        total_rows = metadata.get('total_rows', 0)
        summary_data = [
            ["Metric", "Value"],
            ["Total Records", f"{total_rows:,}"],
            ["Total Columns", str(metadata.get('total_columns', 'N/A'))],
            ["Files Processed", ", ".join(metadata.get('files_processed', ['N/A']))],
        ]
        
        if metadata.get("date_range"):
            summary_data.append(["Date Range", f"{metadata['date_range']['start']} to {metadata['date_range']['end']}"])
        
        table = Table(summary_data, colWidths=[2.5*inch, 4*inch])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), self.primary_color),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), self.light_color),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('GRID', (0, 0), (-1, -1), 1, colors.white),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('TOPPADDING', (0, 1), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 8),
        ]))
        story.append(table)
        
        # Footer
        story.append(Spacer(1, 1*inch))
        story.append(Paragraph(
            f"Report generated by Automated Insight Engine | {insights.get('generated_by', 'AI Analysis')}",
            ParagraphStyle('Footer', parent=styles['Normal'], fontSize=8, 
                          textColor=colors.gray, alignment=TA_CENTER)
        ))
        
        # Build PDF
        doc.build(story)
        return output_path
    
    def _generate_pptx(
        self,
        data: Dict[str, Any],
        insights: Dict[str, Any],
        charts: Dict[str, str],
        config: Dict[str, Any],
        job_id: str
    ) -> str:
        """Generate a professional PowerPoint presentation"""
        output_path = os.path.join(self.output_folder, f"{job_id}.pptx")
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Title Slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = config.get("title", "Performance Report")
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor(37, 99, 235)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add company name
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.75))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = config.get("company_name", "Company")
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RgbColor(59, 130, 246)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Add date
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(0.5))
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = datetime.now().strftime('%B %d, %Y')
        date_para.font.size = Pt(16)
        date_para.font.color.rgb = RgbColor(107, 114, 128)
        date_para.alignment = PP_ALIGN.CENTER
        
        # Slide 2: Executive Summary
        if config.get("include_summary", True) and insights.get("executive_summary"):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_slide_title(slide, "Executive Summary")
            
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            para = tf.paragraphs[0]
            para.text = insights["executive_summary"]
            para.font.size = Pt(16)
            para.font.color.rgb = RgbColor(31, 41, 55)
            para.alignment = PP_ALIGN.LEFT
        
        # Slide 3: Key Findings
        if insights.get("key_findings"):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_slide_title(slide, "Key Findings")
            
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, finding in enumerate(insights["key_findings"][:7]):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.text = f"• {finding}"
                para.font.size = Pt(16)
                para.font.color.rgb = RgbColor(31, 41, 55)
                para.space_before = Pt(10)
        
        # Slide 4-N: Charts
        if config.get("include_charts", True) and charts:
            for chart_name, chart_path in charts.items():
                if os.path.exists(chart_path):
                    try:
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        self._add_slide_title(slide, chart_name.replace("_", " ").title())
                        
                        # Add chart image
                        slide.shapes.add_picture(
                            chart_path,
                            Inches(1.5), Inches(1.75),
                            width=Inches(10.333), height=Inches(5.25)
                        )
                    except Exception as e:
                        print(f"Error adding chart slide {chart_name}: {e}")
        
        # Slide: Trends
        if insights.get("trends"):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_slide_title(slide, "Trends & Patterns")
            
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, trend in enumerate(insights["trends"]):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.text = f"📈 {trend}"
                para.font.size = Pt(18)
                para.font.color.rgb = RgbColor(31, 41, 55)
                para.space_before = Pt(15)
        
        # Slide: Recommendations
        if config.get("include_recommendations", True) and insights.get("recommendations"):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_slide_title(slide, "Strategic Recommendations")
            
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, rec in enumerate(insights["recommendations"][:7]):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.text = f"{i+1}. {rec}"
                para.font.size = Pt(15)
                para.font.color.rgb = RgbColor(31, 41, 55)
                para.space_before = Pt(10)
        
        # Slide: Opportunities
        if insights.get("opportunities"):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_slide_title(slide, "Growth Opportunities")
            
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, opp in enumerate(insights["opportunities"]):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.text = f"💡 {opp}"
                para.font.size = Pt(18)
                para.font.color.rgb = RgbColor(31, 41, 55)
                para.space_before = Pt(12)
        
        # Final Slide: Thank You
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        thank_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1))
        tf = thank_box.text_frame
        para = tf.paragraphs[0]
        para.text = "Thank You"
        para.font.size = Pt(44)
        para.font.bold = True
        para.font.color.rgb = RgbColor(37, 99, 235)
        para.alignment = PP_ALIGN.CENTER
        
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.5))
        tf = footer_box.text_frame
        para = tf.paragraphs[0]
        para.text = f"Powered by Automated Insight Engine | {insights.get('generated_by', 'AI Analysis')}"
        para.font.size = Pt(14)
        para.font.color.rgb = RgbColor(107, 114, 128)
        para.alignment = PP_ALIGN.CENTER
        
        # Save presentation
        prs.save(output_path)
        return output_path
    
    def _add_slide_title(self, slide, title_text: str):
        """Add a consistent title to a slide"""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        para = tf.paragraphs[0]
        para.text = title_text
        para.font.size = Pt(32)
        para.font.bold = True
        para.font.color.rgb = RgbColor(37, 99, 235)
